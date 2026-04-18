"""
parsers.py – Extract text content from PowerPoint and Excel files.
"""
from __future__ import annotations

import io
from typing import List

from pptx import Presentation
import openpyxl


def parse_pptx(file_bytes: bytes) -> List[dict]:
    """
    Parse a PowerPoint file and return a list of slide dicts:
        {"source": "Slide N", "text": "<extracted text>"}
    """
    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
            # Tables inside slides
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        texts.append(" | ".join(cells))
        if texts:
            slides.append({"source": f"Slide {idx}", "text": "\n".join(texts)})
    return slides


def parse_xlsx(file_bytes: bytes) -> List[dict]:
    """
    Parse an Excel file and return a list of sheet dicts:
        {"source": "Sheet: <name>", "text": "<extracted text>"}
    """
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_text = []
        # Collect header row separately
        header = None
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            non_empty = [c for c in cells if c]
            if not non_empty:
                continue
            if header is None:
                header = cells
                rows_text.append(" | ".join(cells))
            else:
                rows_text.append(" | ".join(cells))
        if rows_text:
            sheets.append({"source": f"Sheet: {sheet_name}", "text": "\n".join(rows_text)})
    return sheets
